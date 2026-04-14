# generate_ppt.py - 将所有图表插入PPT
# DeepKsite 项目

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    figures_dir = os.path.join(base_dir, '..', 'figures')
    scatter_dir = os.path.join(figures_dir, 'scatter')
    bar_dir = os.path.join(figures_dir, 'bar')
    output_path = os.path.join(figures_dir, 'DeepKsite_figures.pptx')

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ===== 封面页 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "DeepKsite"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Scatter Plots & Bar Charts"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

    # ===== 散点图 - 测试集（6张，每页2张）=====
    test_scatter_files = [
        ('scatter_test_AAC.png', 'AAC - Test Set'),
        ('scatter_test_DPC.png', 'DPC - Test Set'),
        ('scatter_test_PAAC.png', 'PAAC - Test Set'),
        ('scatter_test_CKSAAGP.png', 'CKSAAGP - Test Set'),
        ('scatter_test_PHYC.png', 'PHYC - Test Set'),
        ('scatter_test_DeepKsite.png', 'DeepKsite - Test Set'),
    ]

    train_scatter_files = [
        ('scatter_train_AAC.png', 'AAC - Train Set'),
        ('scatter_train_DPC.png', 'DPC - Train Set'),
        ('scatter_train_PAAC.png', 'PAAC - Train Set'),
        ('scatter_train_CKSAAGP.png', 'CKSAAGP - Train Set'),
        ('scatter_train_PHYC.png', 'PHYC - Train Set'),
        ('scatter_train_DeepKsite.png', 'DeepKsite - Train Set'),
    ]

    bar_files = [
        ('bar_Logistic_Regression.png', 'Logistic Regression'),
        ('bar_KNN.png', 'KNN'),
        ('bar_Random_Forest.png', 'Random Forest'),
        ('bar_SVM.png', 'SVM'),
        ('bar_Naive_Bayes.png', 'Naive Bayes'),
        ('bar_Gradient_Boosting.png', 'Gradient Boosting'),
    ]

    def add_title_slide(prs, title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        return slide

    def add_two_images_slide(prs, img1_path, img2_path, title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 标题
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        # 左图
        if os.path.exists(img1_path):
            slide.shapes.add_picture(img1_path, Inches(0.3), Inches(1.0), Inches(6.2), Inches(5.5))
        # 右图
        if os.path.exists(img2_path):
            slide.shapes.add_picture(img2_path, Inches(6.8), Inches(1.0), Inches(6.2), Inches(5.5))
        return slide

    def add_single_image_slide(prs, img_path, title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.0), Inches(10), Inches(5.8))
        return slide

    # ===== 散点图页面 =====
    add_title_slide(prs, "Scatter Plots - Test Set")
    for i in range(0, len(test_scatter_files), 2):
        f1, t1 = test_scatter_files[i]
        f2, t2 = test_scatter_files[i+1]
        add_two_images_slide(prs,
                             os.path.join(scatter_dir, f1),
                             os.path.join(scatter_dir, f2),
                             f"{t1}  |  {t2}")

    add_title_slide(prs, "Scatter Plots - Train Set")
    for i in range(0, len(train_scatter_files), 2):
        f1, t1 = train_scatter_files[i]
        f2, t2 = train_scatter_files[i+1]
        add_two_images_slide(prs,
                             os.path.join(scatter_dir, f1),
                             os.path.join(scatter_dir, f2),
                             f"{t1}  |  {t2}")

    # ===== 柱状图页面 =====
    add_title_slide(prs, "Bar Charts - Classifier Comparison")
    for fname, title in bar_files:
        add_single_image_slide(prs,
                               os.path.join(bar_dir, fname),
                               title)

    # ===== 表格页面 =====
    tables_dir = os.path.join(figures_dir, 'tables')
    table_files = [
        ('table1_sota_comparison.png', 'Table 1. Performance Comparison on Test Dataset'),
        ('table2_ablation_test.png', 'Table 2. Ablation Study on Test Dataset'),
        ('table3_fusion_test.png', 'Table 3. Fusion Comparison on Test Dataset'),
        ('table4_ablation_5fold.png', 'Table 4. Ablation Study (5-Fold Cross Validation)'),
        ('table5_fusion_5fold.png', 'Table 5. Fusion Comparison (5-Fold Cross Validation)'),
    ]

    add_title_slide(prs, "Tables")
    for fname, title in table_files:
        add_single_image_slide(prs,
                               os.path.join(tables_dir, fname),
                               title)

    # ===== 保存 =====
    prs.save(output_path)
    print(f"\n PPT已生成: {output_path}")
    print(f"  共 {len(prs.slides)} 页")


if __name__ == '__main__':
    main()
